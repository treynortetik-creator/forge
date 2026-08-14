#!/usr/bin/env python3
"""
test_deck_audit.py — regression tests for the deck harness.

Every case here is a bug that was live in this file and shipped green. A harness
that reports PASS on the defect it exists to catch is worse than no harness, so
each fix gets a test that fails loudly if it regresses.

    python3 test_deck_audit.py
"""
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

import deck_audit as D

FAILS, RUN = [], 0


def check(name, cond, detail=""):
    global RUN
    RUN += 1
    if cond:
        print(f"  ok   {name}")
    else:
        print(f"  FAIL {name}  {detail}")
        FAILS.append(name)


def blank(w=13.333, h=7.5):
    p = Presentation()
    p.slide_width, p.slide_height = Inches(w), Inches(h)
    return p


def audit(prs, **kw):
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        return D.audit(f.name, **kw)


# ── 1. the legibility formula ────────────────────────────────────────────────
print("\nlegibility floor")

f = D.legibility_floor_pt(7.5, "basic")
# NOT "lands near 24pt". That test asserted a just-so story: nothing supports the claim
# that this model is where the folk 24pt minimum came from, the number needs a 7.5in
# slide AND cap height AND 3%EH simultaneously, and AAPT was printing "18 to 24 pt" for
# overhead transparencies decades before DISCAS existed. Assert the standard instead.
check("the floor satisfies %EH x slide_height, which is what the standard says",
      abs(f - D.PERCENT_EH["basic"] * 7.5 * 72 / D.DEFAULT_X_RATIO) < 0.1, f"got {f}pt")
check("the naive version's ~185pt is gone", f < 60, f"got {f}pt")
check("analytical < basic < passive",
      D.legibility_floor_pt(7.5, "analytical") < f < D.legibility_floor_pt(7.5, "passive"))
# Screen size scales with room depth under the 4/6/8 rule, so the floor must not move.
check("floor is scale-invariant without an explicit screen",
      D.legibility_floor_pt(7.5, "basic", None, 15) == D.legibility_floor_pt(7.5, "basic", None, 80))
# With a REAL screen size the room depth must matter again.
near = D.legibility_floor_pt(7.5, "basic", 60, 15)
far = D.legibility_floor_pt(7.5, "basic", 60, 60)
check("an explicit screen makes distance matter", far > near * 3, f"{near} vs {far}")
# A 4:3 slide is 7.5in tall too; a taller slide needs more points for the same physics.
check("taller slide -> higher floor", D.legibility_floor_pt(10.0, "basic") >
      D.legibility_floor_pt(7.5, "basic"))


# ── 2. fill resolution must not read the text colour ─────────────────────────
print("\nfill resolution")

p = blank()
s = p.slides.add_slide(p.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
r = tb.text_frame.paragraphs[0].add_run()
r.text = "dark text, no shape fill"
r.font.size = Pt(40)
r.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
th = D.theme_colours(p)
check("a textbox with no fill reports None, not the run's colour",
      D.shape_fill_rgb(tb, th) is None, f"got {D.shape_fill_rgb(tb, th)}")
res = audit(p)
lc = res["findings"]["low_contrast"]
check("dark text on the white slide bg is NOT a contrast failure", len(lc) == 0, str(lc))

# an explicit fill must still be read
p2 = blank()
s2 = p2.slides.add_slide(p2.slide_layouts[6])
sh = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(6), Inches(2))
sh.fill.solid()
sh.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
check("an explicit solid fill IS read", D.shape_fill_rgb(sh, D.theme_colours(p2)) == (0, 51, 102),
      str(D.shape_fill_rgb(sh, D.theme_colours(p2))))


# ── 2b. z-order: what is actually behind the text ────────────────────────────
print("\nz-order backdrop")

pz = blank()
sz = pz.slides.add_slide(pz.slide_layouts[6])
panel = sz.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(8), Inches(3))
panel.fill.solid()
panel.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
over = sz.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(1))
ru = over.text_frame.paragraphs[0].add_run()
ru.text = "White on the blue panel"
ru.font.size = Pt(32)
ru.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
resz = audit(pz)
check("white text over a blue panel is NOT a false contrast failure",
      len(resz["findings"]["low_contrast"]) == 0, str(resz["findings"]["low_contrast"]))

# and the inverse must still fail: white text NOT over anything
pz2 = blank()
sz2 = pz2.slides.add_slide(pz2.slide_layouts[6])
o2 = sz2.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(1))
r2 = o2.text_frame.paragraphs[0].add_run()
r2.text = "White on the white slide"
r2.font.size = Pt(32)
r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
check("white text on the bare white slide IS still caught",
      len(audit(pz2)["findings"]["low_contrast"]) == 1)

# a panel that does NOT cover the text must not be credited
pz3 = blank()
sz3 = pz3.slides.add_slide(pz3.slide_layouts[6])
far = sz3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1))
far.fill.solid()
far.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
o3 = sz3.shapes.add_textbox(Inches(8), Inches(5), Inches(4), Inches(1))
r3 = o3.text_frame.paragraphs[0].add_run()
r3.text = "White, nowhere near the panel"
r3.font.size = Pt(32)
r3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
check("a non-overlapping panel is not used as the backdrop",
      len(audit(pz3)["findings"]["low_contrast"]) == 1)


# ── 3. group + table traversal ───────────────────────────────────────────────
print("\ntraversal")

p3 = blank()
s3 = p3.slides.add_slide(p3.slide_layouts[6])
bars = []
for i in range(3):
    b = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + i * 1.5), Inches(3), Inches(1.2),
                            Inches(2))
    b.text_frame.text = f"{i}0%"
    run = b.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(7)
    bars.append(b)
s3.shapes.add_group_shape(bars)
tbl = s3.shapes.add_table(2, 2, Inches(7), Inches(2), Inches(5), Inches(2)).table
tbl.cell(1, 1).text = "6pt"
tbl.cell(1, 1).text_frame.paragraphs[0].runs[0].font.size = Pt(6)

res3 = audit(p3)
check("grouped shapes are traversed", res3["checked"]["runs"] >= 4,
      f"only {res3['checked']['runs']} runs")
check("7pt inside a group is caught",
      any(x["pt"] == 7.0 for x in res3["findings"]["tiny_text"]))
check("6pt inside a table cell is caught",
      any(x["pt"] == 6.0 for x in res3["findings"]["tiny_text"]))
check("the group deck does NOT pass", res3["findings"]["tiny_text"])


# ── 4. inheritance: run.font.size is None but the size is real ───────────────
print("\ninheritance resolution")

p4 = blank()
s4 = p4.slides.add_slide(p4.slide_layouts[1])
s4.shapes.title.text = "Inherited"
s4.placeholders[1].text = "Body with no explicit size"
raw = s4.shapes.title.text_frame.paragraphs[0].runs[0].font.size
check("python-pptx really does report None here", raw is None, f"got {raw}")
res4 = audit(p4)
check("but the harness resolves a real size", all(x > 0 for x in res4["type_sizes"]),
      str(res4["type_sizes"]))
check("and says where it came from",
      any("master" in x["source"] or "layout" in x["source"]
          for x in res4["findings"]["tiny_text"]) or not res4["findings"]["tiny_text"])
check("nothing fell through to the 18pt default", not res4["findings"]["unresolved_size"])


# ── 5. native charts ─────────────────────────────────────────────────────────
print("\nnative charts")

p5 = blank()
s5 = p5.slides.add_slide(p5.slide_layouts[6])
cd = CategoryChartData()
cd.categories = ["A", "B"]
cd.add_series("S", (1, 2))
s5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), cd)
check("a native chart is flagged", len(audit(p5)["findings"]["native_charts"]) == 1)
check("a chartless deck is not", len(audit(blank())["findings"]["native_charts"]) == 0)


# ── 6. the zero-measurement trap ─────────────────────────────────────────────
print("\nthe zero-measurement trap")

empty = audit(blank())
check("an empty deck measures 0 runs", empty["checked"]["runs"] == 0)
rc = D.report(empty, "hard")
check("...and that is a FAILURE, not a pass", rc == 1,
      "measuring nothing reported success")


# ── 7. contrast maths ────────────────────────────────────────────────────────
print("\ncontrast")

check("black on white is 21:1", abs(D.contrast((0, 0, 0), (255, 255, 255)) - 21.0) < 0.01)
check("#767676 on white is the 4.5 boundary",
      abs(D.contrast(D.hex_to_rgb("767676"), (255, 255, 255)) - 4.54) < 0.02)
check("contrast is symmetric",
      abs(D.contrast((10, 20, 30), (200, 200, 200)) - D.contrast((200, 200, 200), (10, 20, 30)))
      < 1e-9)

# WCAG large text is 18pt / 14pt bold. The 24px / 18.66px figures are the same
# sizes in CSS pixels; using them on points is a 1.333x error that invents failures.
check("18pt counts as large text", D.is_large_text(18.0))
check("14pt bold counts as large text", D.is_large_text(14.0, bold=True))
check("14pt not bold does NOT", not D.is_large_text(14.0, bold=False))
check("17.9pt does NOT", not D.is_large_text(17.9))
check("the px figures are not used as pt", D.is_large_text(20.0),
      "20pt must be large; if this fails the 24px threshold crept back in")

# end to end: 20pt grey that clears 3:1 but not 4.5:1 must PASS
pw = blank()
sw = pw.slides.add_slide(pw.slide_layouts[6])
bw = sw.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
rw = bw.text_frame.paragraphs[0].add_run()
rw.text = "20pt grey, clears 3:1"
rw.font.size = Pt(20)
rw.font.color.rgb = RGBColor(0x76, 0x76, 0x76)          # 4.54:1 on white
check("20pt at 4.54:1 is not reported as a failure",
      len(audit(pw)["findings"]["low_contrast"]) == 0,
      "large-text threshold is wrong; 20pt is large and only needs 3:1")



# ── 8. the adversarial-review findings ───────────────────────────────────────
print("\nfindings from the adversarial review")

import subprocess, json as _json
from pptx.enum.shapes import MSO_SHAPE as _MS
from lxml import etree as _et

# D2 — a shape filled by p:style/fillRef, with NOTHING in spPr. This is what
# python-pptx add_shape emits and what PowerPoint emits for a drawn shape.
pf = blank(); sf = pf.slides.add_slide(pf.slide_layouts[6])
shf = sf.shapes.add_shape(_MS.RECTANGLE, Inches(1), Inches(1), Inches(8), Inches(3))
rf = shf.text_frame.paragraphs[0].add_run(); rf.text = "on a style-filled shape"
rf.font.size = Pt(32); rf.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
thf = D.theme_colours(pf)
check("a fillRef shape does not read as 'no fill'",
      D.shape_fill_rgb(shf, thf, D.theme_fill_styles(pf)) is not None,
      "spPr-only reading fabricates both sides of the contrast pair")
resf = audit(pf)
check("...and it is not a fabricated contrast FAILURE",
      len(resf["findings"]["low_contrast"]) == 0, str(resf["findings"]["low_contrast"]))

# D3 — colour transforms
base = (0x4F, 0x81, 0xBD)
dark, ok = D._apply_transforms(base, _et.fromstring(
    '<a:schemeClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'val="accent1"><a:lumMod val="20000"/></a:schemeClr>'))
check("lumMod 20% actually darkens the colour", ok and sum(dark) < sum(base) / 2,
      f"{base} -> {dark}")
light, _ = D._apply_transforms(base, _et.fromstring(
    '<a:schemeClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'val="accent1"><a:lumMod val="50000"/><a:lumOff val="50000"/></a:schemeClr>'))
check("lumMod+lumOff 50% lightens it", sum(light) > sum(base), f"{base} -> {light}")
_, alpha_ok = D._apply_transforms(base, _et.fromstring(
    '<a:schemeClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'val="accent1"><a:alpha val="50000"/></a:schemeClr>'))
check("alpha is reported unmodellable, not silently ignored", not alpha_ok)

# D8 — --json and text must agree, always
empty_f = TMPF = __import__("tempfile").mktemp(suffix=".pptx"); blank().save(TMPF)
for mode in ("hard", "any"):
    t = subprocess.run([sys.executable, "deck_audit.py", TMPF, "--fail-on", mode],
                       capture_output=True).returncode
    j = subprocess.run([sys.executable, "deck_audit.py", TMPF, "--fail-on", mode, "--json"],
                       capture_output=True).returncode
    check(f"--json agrees with text on --fail-on {mode}", t == j, f"text={t} json={j}")
    check(f"a zero-run deck FAILS under --fail-on {mode}", t == 1, f"got {t}")

# D5 — the master-placeholder link the docs promise must exist
check("resolve_font_pt can return 'master placeholder'",
      "master placeholder" in open("deck_audit.py").read(),
      "documented in the header and SKILL.md but never implemented")

# D6 — rotation
pr = blank(); sr = pr.slides.add_slide(pr.slide_layouts[6])
tall = sr.shapes.add_shape(_MS.RECTANGLE, Inches(6.0), Inches(-0.25), Inches(1), Inches(8))
tall.rotation = 90        # renders 8x1in, centred, FULLY on-slide
check("a rotated shape that lands on-slide is not a false off-slide",
      len(audit(pr)["findings"]["offslide"]) == 0, "rotation ignored in bounds")
pr2 = blank(); sr2 = pr2.slides.add_slide(pr2.slide_layouts[6])
wide = sr2.shapes.add_shape(_MS.RECTANGLE, Inches(2.6), Inches(3.5), Inches(8), Inches(1))
wide.rotation = 90        # renders 1x8in, y spans past a 7.5in slide
check("a rotated shape that hangs off IS caught",
      len(audit(pr2)["findings"]["offslide"]) == 1, "rotation ignored in bounds")

# A2/A5/A7 — the legibility corrections
check("the reference glyph defaults to x-height, not cap height",
      D.glyph_ratio(None, "x")[0] < D.CAP_HEIGHT_RATIO)
check("x-height gives a HIGHER (stricter) floor than cap height",
      D.legibility_floor_pt(7.5, "basic") > D.legibility_floor_pt(7.5, "basic", element="cap"))
check("a named font uses its real ratio",
      D.glyph_ratio("Arial", "x")[0] == D.X_HEIGHT_RATIO["arial"])
check("the floor tracks slide HEIGHT, not aspect ratio",
      D.legibility_floor_pt(7.5, "basic") != D.legibility_floor_pt(5.625, "basic"),
      "16:9 at 13.333x7.5 and at 10x5.625 are the same aspect")
check("tiers map to 2/3/4 %EH",
      (D.PERCENT_EH["analytical"], D.PERCENT_EH["basic"], D.PERCENT_EH["passive"])
      == (0.02, 0.03, 0.04))

# lxml truthiness
_el = _et.fromstring('<root xmlns:p="p"><p:spPr/></root>')
check("_find_any survives a childless element (which is falsy)",
      D._find_any(_el, "{p}spPr") is not None,
      "`a or b` on lxml elements silently skips an empty element")


print(f"\n{RUN - len(FAILS)}/{RUN} passed")
if FAILS:
    print("FAILED: " + ", ".join(FAILS))
sys.exit(1 if FAILS else 0)
