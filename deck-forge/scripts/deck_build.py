#!/usr/bin/env python3
"""
deck_build.py — build decks that survive a Google Slides import, by construction.

WHY THIS EXISTS
---------------
A native PowerPoint chart **flattens into a static image** when the file is imported
into Google Slides, and there is no workaround inside an uploaded `.pptx`. Slides
only keeps a chart live when it is linked to a Sheet, and that link cannot exist in
a file you upload. So the chart arrives as a dead picture: not editable, not
restylable, not even selectable as data.

The usual response is a style guide telling people not to use charts. Style guides
lose. So this library **cannot express a native chart** — there is no `add_chart`,
and there is no way to reach one through the public API here. Everything visual is
built from rectangles, ovals, text boxes and tables, which are the primitives that
survive the import editable. Compliance is structural, not advisory.

Two more properties, both deliberate:

- **Every run gets an explicit size and colour.** Inherited formatting is why
  `run.font.size` returns None on a normal deck and why a naive checker reports a
  clean file. Setting both explicitly means `deck_audit.py` can actually measure
  what it built.
- **The legibility floor is enforced at build time**, not reported afterwards. Ask
  for 8pt text and you get an exception, not a warning you will scroll past.

THE HONEST TRADE-OFF, which you should state to the user once: shape bars are not
data-driven. Changing a number means resizing a rectangle and editing its label,
which this library does for you on a rebuild but Slides will not do on its own.
That is the price of surviving the import editable.

USAGE
    from deck_build import Deck, Theme
    d = Deck(Theme.load("theme.json"))
    s = d.slide("Title of this slide")
    s.bars([("Q1", 42), ("Q2", 61), ("Q3", 88)], unit="%")
    d.save("out.pptx")

    python3 deck_build.py --demo out.pptx     # build a demo and audit it

Requires python-pptx.
"""
import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("ERROR: needs python-pptx.  pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)

sys.path.insert(0, str(Path(__file__).parent))
try:
    from deck_audit import legibility_floor_pt
except ImportError:                                    # standalone use
    def legibility_floor_pt(slide_h_in, need="basic", *a, **k):
        return round({"analytical": .02, "basic": .03, "passive": .04}[need]
                     * slide_h_in * 72.0 / 0.52, 1)


DEFAULT_THEME = {
    "name": "Neutral",
    "slide": {"width_in": 13.333, "height_in": 7.5},
    "colours": {
        "bg": "FFFFFF", "ink": "1A1A1A", "muted": "5A5A5A",
        "accent": "1F4E79", "accent2": "C55A11", "rule": "D8D8D8",
        "on_accent": "FFFFFF",
    },
    # Every size clears the 'basic' floor, which is 31.2pt on a 7.5in slide once the
    # reference glyph is the LOWERCASE letter AVIXA's own training specifies rather
    # than the capital this library first used. That correction moved the floor up by
    # ~35%, and these sizes moved with it. They look large because most decks are in
    # fact too small for the room they are shown in; that is the finding, not a bug.
    # If you need denser type you are authoring a document to be read up close — set
    # viewing_need to "analytical" (2%EH), and only if the screen really is that big
    # relative to the room.
    "type": {"family": "Arial", "title_pt": 44, "subtitle_pt": 32,
             "body_pt": 32, "label_pt": 32, "stat_pt": 66},
    "layout": {"margin_in": 0.9, "gutter_in": 0.35, "title_top_in": 0.6},
    "viewing_need": "basic",
}


MIN_SHAPE_IN = 0.02          # below this a shape is a hairline or schema-invalid
VIEWING_NEEDS = ("analytical", "basic", "passive")
HEX = __import__("re").compile(r"^#?[0-9A-Fa-f]{6}$")


class DeckError(Exception):
    """Raised when a deck would be built wrong. Loud on purpose."""


def _lum(rgb):
    def f(c):
        c /= 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * f(rgb[0]) + 0.7152 * f(rgb[1]) + 0.0722 * f(rgb[2])


def contrast_ratio(a, b):
    l1, l2 = sorted((_lum(a), _lum(b)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


class Theme:
    def __init__(self, data=None):
        self.d = json.loads(json.dumps(DEFAULT_THEME))
        if data:
            for k, v in data.items():
                if isinstance(v, dict) and isinstance(self.d.get(k), dict):
                    self.d[k].update(v)
                else:
                    self.d[k] = v
        # The deep copy above means every default key is always present, so the old
        # "missing colours" guard could never fire for a dict. Validate the VALUES.
        for k, v in self.d["colours"].items():
            if not isinstance(v, str) or not HEX.match(v):
                raise DeckError(
                    f"theme colour {k!r} is {v!r}; expected a 6-digit hex like "
                    f"'1F4E79' or '#1F4E79'. Named CSS colours and 3-digit hex are "
                    f"not supported.")
        for k, v in self.d["type"].items():
            if k == "family":
                if not isinstance(v, str) or not v.strip():
                    raise DeckError(f"theme type.family is {v!r}; expected a font name")
            elif not isinstance(v, (int, float)) or v <= 0:
                raise DeckError(f"theme type.{k} is {v!r}; expected a positive number")
        for k, v in self.d["layout"].items():
            if not isinstance(v, (int, float)) or v < 0:
                raise DeckError(f"theme layout.{k} is {v!r}; expected a number >= 0")
        sl = self.d["slide"]
        for k in ("width_in", "height_in"):
            if not isinstance(sl.get(k), (int, float)) or sl[k] <= 0:
                raise DeckError(f"theme slide.{k} is {sl.get(k)!r}; expected a positive number")

    @classmethod
    def load(cls, path):
        try:
            data = json.loads(Path(path).read_text("utf-8"))
        except FileNotFoundError:
            raise DeckError(f"theme file not found: {path}")
        except json.JSONDecodeError as e:
            raise DeckError(f"theme file {path} is not valid JSON: {e}")
        if not isinstance(data, dict):
            raise DeckError(f"theme file {path} must contain a JSON object")
        try:
            return cls(data)
        except DeckError as e:
            raise DeckError(f"{path}: {e}")

    def rgb(self, name):
        v = self.d["colours"].get(name, name)
        return RGBColor.from_string(str(v).lstrip("#").upper())

    def pt(self, name):
        return self.d["type"][name]

    @property
    def family(self):
        return self.d["type"]["family"]


class Slide:
    """One slide. Every method here builds from primitives that survive the import."""

    def __init__(self, deck, native, title=None, subtitle=None):
        self.deck, self.s, self.t = deck, native, deck.theme
        m = self.t.d["layout"]["margin_in"]
        self.left, self.right = m, self.deck.w - m
        self.y = self.t.d["layout"]["title_top_in"]
        if title:
            self.heading(title, subtitle)

    # ── primitives ───────────────────────────────────────────────────────────

    def _text(self, x, y, w, h, text, pt, colour="ink", bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, on="bg"):
        self.deck._check_pt(pt, text)
        self.deck._check_contrast(colour, on, pt, bold, text)
        box = self.s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        # Explicit, always. Inherited formatting is exactly what makes a deck
        # unmeasurable — and unmeasured is the state this whole plugin exists to end.
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.name = self.t.family
        run.font.color.rgb = self.t.rgb(colour)
        return box

    def _rect(self, x, y, w, h, fill="accent", line=None, shape=MSO_SHAPE.RECTANGLE):
        sh = self.s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self.t.rgb(fill)
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = self.t.rgb(line)
        sh.shadow.inherit = False
        if sh.has_text_frame:
            sh.text_frame.word_wrap = True
        return sh

    # ── components ───────────────────────────────────────────────────────────

    def heading(self, title, subtitle=None):
        w = self.right - self.left
        self._text(self.left, self.y, w, 1.0, title, self.t.pt("title_pt"), "ink", bold=True)
        self.y += 1.05
        if subtitle:
            self._text(self.left, self.y, w, 0.7, subtitle, self.t.pt("subtitle_pt"), "muted")
            self.y += 0.75
        self._rect(self.left, self.y, w, 0.02, fill="rule")
        self.y += 0.45
        return self

    def bullets(self, items, pt=None):
        pt = pt or self.t.pt("body_pt")
        w = self.right - self.left
        for item in items:
            self._rect(self.left, self.y + pt / 72 * 0.38, 0.13, 0.13, fill="accent",
                       shape=MSO_SHAPE.OVAL)
            self._text(self.left + 0.35, self.y, w - 0.35, pt / 72 * 1.6, item, pt, "ink")
            self.y += pt / 72 * 1.75
            self.deck._check_bounds(self.y, 'bullets')
        return self

    def bars(self, data, unit="", height_in=3.2, colour="accent"):
        """A column chart made of rectangles. Editable in Slides; a real chart is not."""
        if not data:
            raise DeckError("bars() needs at least one (label, value) pair")
        vals = [float(v) for _, v in data]
        if min(vals) < 0:
            raise DeckError("bars() does not do negative values; use a different component "
                            "rather than a bar that points the wrong way")
        top = max(vals) or 1.0
        n = len(data)
        gut = self.t.d["layout"]["gutter_in"]
        total_w = self.right - self.left
        bw = (total_w - gut * (n - 1)) / n
        # a:ext/@cx is ST_PositiveCoordinate (minInclusive 0). Past ~30 bars this goes
        # negative and emits schema-invalid OOXML that PowerPoint may refuse to open --
        # and the audit certified 150 negative-width shapes as clean, exit 0.
        if bw < MIN_SHAPE_IN:
            raise DeckError(
                f"{n} bars do not fit across {total_w:.2f}in: each would be "
                f"{bw:.3f}in wide, and a shape cannot have negative or hairline width "
                f"(a:ext/@cx must be >= 0).\nAt most "
                f"{int((total_w + gut) // (MIN_SHAPE_IN + gut))} bars fit here. Split the "
                f"data across slides, or use a table.")
        label_pt = self.t.pt("label_pt")
        base = self.y + height_in
        for i, (label, val) in enumerate(data):
            x = self.left + i * (bw + gut)
            bh = max(float(val) / top * (height_in - 0.7), 0.04)
            self._rect(x, base - bh, bw, bh, fill=colour)
            self._text(x, base - bh - 0.5, bw, 0.45, f"{val}{unit}", label_pt, "ink",
                       bold=True, align=PP_ALIGN.CENTER)
            self._text(x, base + 0.1, bw, 0.5, label, label_pt, "muted", align=PP_ALIGN.CENTER)
        self._rect(self.left, base, total_w, 0.02, fill="rule")
        self.y = base + 0.75
        self.deck._check_bounds(self.y, 'bars')
        return self

    def stats(self, items, colour="accent"):
        """A row of big numbers. items = [(value, caption), ...]"""
        n = len(items)
        if not n:
            raise DeckError("stats() needs at least one item")
        gut = self.t.d["layout"]["gutter_in"]
        total_w = self.right - self.left
        cw = (total_w - gut * (n - 1)) / n
        for i, (value, caption) in enumerate(items):
            x = self.left + i * (cw + gut)
            self._text(x, self.y, cw, 1.1, value, self.t.pt("stat_pt"), colour, bold=True)
            self._text(x, self.y + 1.15, cw, 0.9, caption, self.t.pt("label_pt"), "muted")
        self.y += 2.2
        self.deck._check_bounds(self.y, 'stats')
        return self

    def callout(self, text, attribution=None, fill="accent", ink="on_accent"):
        """Text INSIDE the filled shape, never floating over it.

        A separate text box laid on top of a coloured rectangle looks identical in
        PowerPoint, but the text then has no fill of its own, so its background is a
        z-order question rather than a fact. Putting the runs in the shape's own text
        frame makes the contrast pair real and keeps it one object through the import.
        """
        w = self.right - self.left
        pt = self.t.pt("body_pt")
        lines = max(1, len(str(text)) // 60 + 1)
        h = 0.6 + lines * (pt / 72 * 1.35) + (0.5 if attribution else 0)
        box = self._rect(self.left, self.y, w, h, fill=fill)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.4)
        for i, (val, size) in enumerate([(text, pt)] +
                                        ([(attribution, self.t.pt("label_pt"))]
                                         if attribution else [])):
            self.deck._check_pt(size, val)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(size)
            run.font.name = self.t.family
            run.font.color.rgb = self.t.rgb(ink)
        self.y += h + 0.4
        self.deck._check_bounds(self.y, 'block')
        return self

    def table(self, rows, col_widths=None, header=True):
        if not rows or not rows[0]:
            raise DeckError("table() needs at least one row with one cell")
        nrow, ncol = len(rows), len(rows[0])
        if any(len(r) != ncol for r in rows):
            raise DeckError("every table row must have the same number of cells")
        w = self.right - self.left
        pt = self.t.pt("label_pt")
        h = nrow * (pt / 72 * 2.0)
        gt = self.s.shapes.add_table(nrow, ncol, Inches(self.left), Inches(self.y),
                                     Inches(w), Inches(h)).table
        if col_widths:
            if len(col_widths) != ncol:
                raise DeckError("col_widths must have one entry per column")
            scale = w / sum(col_widths)
            for i, cw in enumerate(col_widths):
                gt.columns[i].width = Inches(cw * scale)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = gt.cell(r, c)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.size = Pt(pt)
                run.font.name = self.t.family
                is_head = header and r == 0
                run.font.bold = is_head
                run.font.color.rgb = self.t.rgb("on_accent" if is_head else "ink")
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.t.rgb("accent" if is_head else "bg")
        self.deck._check_pt(pt, "table text")
        self.y += h + 0.4
        self.deck._check_bounds(self.y, 'block')
        return self


class Deck:
    def __init__(self, theme=None):
        self.theme = theme or Theme()
        self.w = self.theme.d["slide"]["width_in"]
        self.h = self.theme.d["slide"]["height_in"]
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(self.w), Inches(self.h)
        self.need = self.theme.d.get("viewing_need", "basic")
        # .get(need, default) swallowed typos: a theme meaning "passive" (floor 30.9pt)
        # that misspelled it silently got 23.1pt, and text that should have been
        # rejected built and audited clean.
        if self.need not in VIEWING_NEEDS:
            raise DeckError(
                f"unknown viewing_need {self.need!r}. Use one of: "
                f"{', '.join(sorted(VIEWING_NEEDS))}. Guessing here would quietly "
                f"lower the legibility floor, which is the one thing this must not do.")
        self.floor_pt = legibility_floor_pt(self.h, self.need)
        self.slides = []

    def _check_contrast(self, fg_name, bg_name, pt, bold, what=""):
        """Contrast is half this plugin's headline; enforcing it only after the fact
        let a theme with white ink on a white background build without complaint."""
        try:
            fg = tuple(self.theme.rgb(fg_name))
            bg = tuple(self.theme.rgb(bg_name))
        except Exception:
            return
        ratio = contrast_ratio(fg, bg)
        # WCAG large text is 18pt, or 14pt bold. (24px/18.66px are those same sizes
        # in CSS pixels — applying the pixel figures to points is a 1.333x error.)
        need = 3.0 if (pt >= 18.0 or (bold and pt >= 14.0)) else 4.5
        if ratio < need:
            raise DeckError(
                f"{fg_name} on {bg_name} is {ratio:.2f}:1, below the {need}:1 WCAG "
                f"minimum for {pt}pt{' bold' if bold else ''} text"
                + (f" (text: {str(what)[:40]!r})" if what else "")
                + ".\nChange the colour pair in the theme. Do not ship text nobody "
                  "can read.")

    def _check_bounds(self, y, what=""):
        if y > self.h + 0.01:
            raise DeckError(
                f"content reached {y:.2f}in on a {self.h}in slide — {y - self.h:.2f}in "
                f"past the bottom edge" + (f" at {what}" if what else "") +
                ".\nSplit across slides or cut content. Off-slide content is invisible "
                "in the room and silently truncated in a PDF export.")

    def _check_pt(self, pt, what=""):
        if pt < self.floor_pt:
            raise DeckError(
                f"{pt}pt is below the {self.floor_pt}pt legibility floor for a "
                f"'{self.need}' viewing need on a {self.h}in-tall slide"
                + (f" (text: {str(what)[:40]!r})" if what else "")
                + ".\nRaise the size, cut the words, or pass a different viewing_need in "
                  "the theme. Do not ship text nobody in the room can read.")

    def slide(self, title=None, subtitle=None):
        native = self.prs.slides.add_slide(self.prs.slide_layouts[6])   # blank
        bg = self.theme.d["colours"]["bg"]
        if bg.upper().lstrip("#") != "FFFFFF":
            s0 = Slide(self, native)
            s0._rect(0, 0, self.w, self.h, fill="bg")
        s = Slide(self, native, title, subtitle)
        self.slides.append(s)
        return s

    def save(self, path):
        if not self.slides:
            raise DeckError("refusing to save a deck with no slides")
        self.prs.save(str(path))
        return path


# ── demo, which doubles as the builder↔harness round trip ────────────────────

def demo(out):
    d = Deck(Theme())
    s = d.slide("Everything here is a rectangle",
                "No native charts, because they flatten on import")
    s.bars([("Q1", 42), ("Q2", 61), ("Q3", 88), ("Q4", 74)], unit="%")
    s2 = d.slide("Numbers that survive the trip")
    s2.stats([("88%", "of the thing, measured"), ("4.5:1", "minimum contrast"),
              ("23pt", "legibility floor")])
    s2.callout("A verification tool that itself needs verifying is worse than nothing.",
               "the constraint this plugin is built around")
    s3 = d.slide("Tables survive too")
    s3.table([["Element", "Survives import", "Editable after"],
              ["Rectangle", "yes", "yes"],
              ["Text box", "yes", "yes"],
              ["Table", "yes", "yes"],
              ["Native chart", "as a picture", "no"]])
    d.save(out)
    return d


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--demo", metavar="OUT.pptx", help="build a demo deck")
    ap.add_argument("--theme", type=Path, help="theme JSON")
    ap.add_argument("--print-theme", action="store_true", help="dump the default theme")
    a = ap.parse_args()
    if a.print_theme:
        print(json.dumps(DEFAULT_THEME, indent=2))
        return 0
    if a.demo:
        d = demo(a.demo)
        print(f"  built {a.demo} — {len(d.slides)} slides, floor {d.floor_pt}pt")
        print(f"  now audit it:  python3 deck_audit.py {a.demo}")
        return 0
    ap.print_help()
    return 0


if __name__ == "__main__":
    sys.exit(main())
