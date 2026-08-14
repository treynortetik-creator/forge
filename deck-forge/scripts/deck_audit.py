#!/usr/bin/env python3
"""
deck_audit.py — measure a .pptx instead of eyeballing it.

WHY THIS EXISTS
---------------
Same thesis as the page harness: a person looking at a slide cannot tell whether the
18pt caption will be legible from the back row, whether a colour pair clears 4.5:1,
or whether a native chart is about to flatten into a dead image on import. All of
those are arithmetic.

🔴 THE TRAP THIS IS BUILT AROUND. `run.font.size` returns **None** whenever the size
is inherited from a layout, a master, or the master's txStyles — which is the normal
case for placeholder text. A naive reader sees None, skips the shape, and reports a
clean deck. So this resolves the real inheritance chain before measuring anything:

    run.rPr → paragraph.defRPr → layout placeholder → master placeholder
            → master txStyles (title / body / other) → 18pt PowerPoint default

Every check reports `checked: n` and refuses to pass at zero, for the same reason.

USAGE
    python3 deck_audit.py deck.pptx
    python3 deck_audit.py deck.pptx --room-depth 30 --json
    python3 deck_audit.py deck.pptx --fail-on any

Requires python-pptx.
"""
import argparse
import json
import math
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: needs python-pptx.  pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
EMU_PER_IN = 914400


# ── colour ───────────────────────────────────────────────────────────────────

def _srgb(el):
    """Pull an explicit sRGB value out of a fill/colour element, if there is one."""
    if el is None:
        return None
    c = el.find(f"{A}srgbClr")
    if c is not None and c.get("val"):
        return c.get("val").upper()
    return None


def theme_colours(prs):
    """Map scheme names (dk1, lt1, accent1…) to concrete RGB from the theme part."""
    out = {}
    try:
        theme = prs.slide_master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        root = theme._element if hasattr(theme, "_element") else None
        if root is None:
            from lxml import etree
            root = etree.fromstring(theme.blob)
        scheme = root.find(f".//{A}clrScheme")
        if scheme is not None:
            for child in scheme:
                name = child.tag.split("}")[-1]
                val = _srgb(child)
                if val is None:
                    sysc = child.find(f"{A}sysClr")
                    if sysc is not None and sysc.get("lastClr"):
                        val = sysc.get("lastClr").upper()
                if val:
                    out[name] = val
    except Exception:
        pass
    return out


def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def rel_lum(rgb):
    def f(c):
        c /= 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * f(r) + 0.7152 * f(g) + 0.0722 * f(b)


def contrast(a, b):
    l1, l2 = sorted((rel_lum(a), rel_lum(b)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


# WCAG 2.x SC 1.4.3 defines large-scale text as **18 point, or 14 point bold**.
# The familiar 24px / 18.66px figures are those same sizes expressed in CSS pixels
# at 96dpi (18pt = 24px, 14pt = 18.67px). A .pptx is measured in POINTS, so using
# the pixel numbers here is a 1.333x unit error that holds every run between 18pt
# and 24pt to 4.5:1 when the spec asks 3:1 -- manufacturing failures that are not
# failures. This harness had exactly that bug, copied from its browser sibling.
LARGE_PT, LARGE_BOLD_PT = 18.0, 14.0


def is_large_text(pt, bold=False):
    return pt >= LARGE_PT or (bold and pt >= LARGE_BOLD_PT)


# ── the inheritance resolver — the reason this file exists ───────────────────

def _size_from_rpr(rpr):
    if rpr is not None and rpr.get("sz"):
        return int(rpr.get("sz")) / 100.0
    return None


def _ph_idx(shape):
    try:
        pf = shape.placeholder_format
        return pf.idx, str(pf.type)
    except Exception:
        return None, None


def _find_layout_ph(layout, idx):
    for ph in layout.placeholders:
        try:
            if ph.placeholder_format.idx == idx:
                return ph
        except Exception:
            continue
    return None


def _lvl_size_from_liststyle(el, lvl=0):
    """Read lvlNpPr/defRPr@sz out of a lstStyle or txStyles block."""
    if el is None:
        return None
    tag = f"{A}lvl{lvl + 1}pPr"
    node = el.find(tag)
    if node is None:
        return None
    return _size_from_rpr(node.find(f"{A}defRPr"))


def resolve_font_pt(run, para, shape, slide, prs, master_styles):
    """
    Effective point size for a run, walking the real chain.
    Returns (points, source) so the report can say WHERE the size came from.
    """
    # 1. explicit on the run
    if run is not None and run.font.size is not None:
        return run.font.size.pt, "run"
    rpr = run._r.find(f"{A}rPr") if run is not None else None
    v = _size_from_rpr(rpr)
    if v:
        return v, "run"

    # 2. paragraph defaults
    ppr = para._p.find(f"{A}pPr") if para is not None else None
    if ppr is not None:
        v = _size_from_rpr(ppr.find(f"{A}defRPr"))
        if v:
            return v, "paragraph"
    lvl = 0
    if ppr is not None and ppr.get("lvl"):
        lvl = int(ppr.get("lvl"))

    # 3. the shape's own lstStyle
    tx = shape.text_frame._txBody if shape.has_text_frame else None
    if tx is not None:
        v = _lvl_size_from_liststyle(tx.find(f"{A}lstStyle"), lvl)
        if v:
            return v, "shape lstStyle"

    idx, ph_type = _ph_idx(shape)

    # 4. matching placeholder on the layout
    if idx is not None:
        lay_ph = _find_layout_ph(slide.slide_layout, idx)
        if lay_ph is not None and lay_ph.has_text_frame:
            body = lay_ph.text_frame._txBody
            v = _lvl_size_from_liststyle(body.find(f"{A}lstStyle"), lvl)
            if v:
                return v, "layout placeholder"
            for p2 in body.findall(f"{A}p"):
                v = _size_from_rpr(p2.find(f"{A}pPr/{A}defRPr"))
                if v:
                    return v, "layout placeholder"

    # 5. the matching placeholder on the MASTER. Documented in this module's header
    #    and in SKILL.md, but previously not implemented -- so a 6pt size set on a
    #    master placeholder was skipped and master txStyles answered 32pt instead.
    #    Corporate templates routinely size on master placeholders.
    if idx is not None:
        try:
            mst_ph = _find_layout_ph(slide.slide_layout.slide_master, idx)
            if mst_ph is not None and mst_ph.has_text_frame:
                body = mst_ph.text_frame._txBody
                v = _lvl_size_from_liststyle(body.find(f"{A}lstStyle"), lvl)
                if v:
                    return v, "master placeholder"
                for p3 in body.findall(f"{A}p"):
                    v = _size_from_rpr(p3.find(f"{A}pPr/{A}defRPr"))
                    if v:
                        return v, "master placeholder"
        except Exception:
            pass

    # 6. master txStyles, keyed by placeholder role
    key = "body"
    t = (ph_type or "").lower()
    if "title" in t or "ctrtitle" in t:
        key = "title"
    elif idx is None:
        key = "other"
    v = _lvl_size_from_liststyle(master_styles.get(key), lvl)
    if v:
        return v, f"master txStyles/{key}"

    # 6. PowerPoint's own default
    return 18.0, "PowerPoint default (18pt)"


def _colour_from_rpr(rpr, theme):
    if rpr is None:
        return None
    fill = rpr.find(f"{A}solidFill")
    return _solid_to_rgb(fill, theme) if fill is not None else None


def _lvl_colour_from_liststyle(el, theme, lvl=0):
    if el is None:
        return None
    node = el.find(f"{A}lvl{lvl + 1}pPr")
    return _colour_from_rpr(node.find(f"{A}defRPr"), theme) if node is not None else None


def resolve_font_colour(run, para, shape, slide, prs, master_styles, theme):
    """Effective text colour, walking the SAME chain as resolve_font_pt.

    Without this the contrast check only covers runs with an explicit colour, which on
    a real deck is a small minority -- everything else inherits from the layout or the
    master. Reporting those as 'unmeasured' is honest but nearly useless; resolving
    them is what makes the check actually cover the deck.
    """
    got = run_colour_rgb(run, theme)
    if got:
        return got, "run"
    ppr = para._p.find(f"{A}pPr") if para is not None else None
    lvl = int(ppr.get("lvl")) if (ppr is not None and ppr.get("lvl")) else 0
    if ppr is not None:
        got = _colour_from_rpr(ppr.find(f"{A}defRPr"), theme)
        if got:
            return got, "paragraph"
    if shape.has_text_frame:
        got = _lvl_colour_from_liststyle(
            shape.text_frame._txBody.find(f"{A}lstStyle"), theme, lvl)
        if got:
            return got, "shape lstStyle"
    idx, ph_type = _ph_idx(shape)
    if idx is not None:
        lay_ph = _find_layout_ph(slide.slide_layout, idx)
        if lay_ph is not None and lay_ph.has_text_frame:
            got = _lvl_colour_from_liststyle(
                lay_ph.text_frame._txBody.find(f"{A}lstStyle"), theme, lvl)
            if got:
                return got, "layout placeholder"
    key = "body"
    t = (ph_type or "").lower()
    if "title" in t or "ctrtitle" in t:
        key = "title"
    elif idx is None:
        key = "other"
    got = _lvl_colour_from_liststyle(master_styles.get(key), theme, lvl)
    if got:
        return got, f"master txStyles/{key}"
    try:
        ref = shape._element.find(f".//{P}style/{A}fontRef")
        if ref is not None:
            rgb, ok = _colour_el_to_rgb(_first_colour_child(ref), theme)
            if rgb is not None:
                return (rgb, "p:style/fontRef") if ok else (None, "fontRef has an "
                                                            "unmodellable transform")
    except Exception:
        pass
    # 🔴 DO NOT fall back to dk1. Assuming black here is what stopped the
    # unmeasured_contrast channel from EVER firing: the escape hatch existed, and
    # every unresolved run took this branch instead and got a confident verdict
    # against a colour nobody wrote. Unknown must stay unknown.
    return None, "unresolved"


def master_txstyles(prs):
    out = {}
    try:
        tx = prs.slide_master.element.find(f"{P}txStyles")
        if tx is not None:
            for name, tag in (("title", "titleStyle"), ("body", "bodyStyle"), ("other", "otherStyle")):
                out[name] = tx.find(f"{P}{tag}")
    except Exception:
        pass
    return out


# ── shape colour resolution ──────────────────────────────────────────────────

UNKNOWN_FILL = "UNKNOWN"   # picture or gradient: refuse to guess, skip the check

SCHEME_ALIAS = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}


def _find_any(el, *tags):
    """First matching DIRECT child. Never use `a or b` on lxml elements: an element
    with no children is FALSY, so a childless <p:spPr/> silently falls through to the
    next lookup and reports None. Same trap as `if not element:`."""
    for t in tags:
        got = el.find(t)
        if got is not None:
            return got
    return None


def _hsl(r, g, b):
    r, g, b = r / 255, g / 255, b / 255
    mx, mn = max(r, g, b), min(r, g, b)
    l = (mx + mn) / 2
    if mx == mn:
        return 0.0, 0.0, l
    d = mx - mn
    sat = d / (2 - mx - mn) if l > 0.5 else d / (mx + mn)
    if mx == r:
        h = ((g - b) / d + (6 if g < b else 0)) / 6
    elif mx == g:
        h = ((b - r) / d + 2) / 6
    else:
        h = ((r - g) / d + 4) / 6
    return h, sat, l


def _rgb(h, sat, l):
    def hue(p, q, t):
        t %= 1
        if t < 1 / 6: return p + (q - p) * 6 * t
        if t < 1 / 2: return q
        if t < 2 / 3: return p + (q - p) * (2 / 3 - t) * 6
        return p
    if sat == 0:
        v = int(round(l * 255)); return (v, v, v)
    q = l * (1 + sat) if l < 0.5 else l + sat - l * sat
    pp = 2 * l - q
    return tuple(int(round(max(0, min(1, hue(pp, q, h + o))) * 255))
                 for o in (1 / 3, 0, -1 / 3))


def _pct(el, tag):
    n = el.find(f"{A}{tag}")
    if n is None or not n.get("val"):
        return None
    return int(n.get("val")) / 100000.0


def _apply_transforms(rgb, clr_el):
    """Apply lumMod / lumOff / shade / tint to a resolved colour.

    PowerPoint's colour picker emits lumMod+lumOff for EVERY "Lighter/Darker %"
    swatch in the theme row, so a deck styled from the UI is full of them. Reading
    only schemeClr@val returns the BASE colour, which for `accent1 lumMod 20000` is a
    mid blue where the slide actually shows near-black -- reporting 4.31:1 on a pair
    that renders at about 1.2:1, silently.

    Returns (rgb, ok). ok=False means a transform is present that this cannot model,
    in which case the caller must report the colour as indeterminate rather than
    guess. Alpha is the main one: partial transparency needs compositing against
    whatever is behind, which is not a colour question.
    """
    if clr_el is None:
        return rgb, True
    if clr_el.find(f"{A}alpha") is not None:
        return rgb, False
    known = {"lumMod", "lumOff", "shade", "tint", "satMod", "satOff",
             "hueMod", "hueOff", "comp", "inv", "gray", "gamma", "invGamma"}
    for child in clr_el:
        tag = child.tag.split("}")[-1]
        if tag in ("comp", "inv", "gray", "gamma", "invGamma"):
            return rgb, False

    h, sat, l = _hsl(*rgb)
    v = _pct(clr_el, "lumMod")
    if v is not None:
        l *= v
    v = _pct(clr_el, "lumOff")
    if v is not None:
        l += v
    v = _pct(clr_el, "satMod")
    if v is not None:
        sat *= v
    v = _pct(clr_el, "satOff")
    if v is not None:
        sat += v
    v = _pct(clr_el, "hueMod")
    if v is not None:
        h *= v
    l, sat = max(0.0, min(1.0, l)), max(0.0, min(1.0, sat))
    out = _rgb(h, sat, l)

    # shade/tint operate on the linear channels, not on HSL luminance.
    v = _pct(clr_el, "shade")
    if v is not None:
        out = tuple(int(round(c * v)) for c in out)
    v = _pct(clr_el, "tint")
    if v is not None:
        out = tuple(int(round(c * v + 255 * (1 - v))) for c in out)
    return tuple(max(0, min(255, c)) for c in out), True


def _colour_el_to_rgb(el, theme):
    """Any DrawingML colour element (srgbClr / schemeClr / sysClr / prstClr) -> RGB.

    Returns (rgb, ok); ok=False when a transform is present that cannot be modelled.
    """
    if el is None:
        return None, True
    tag = el.tag.split("}")[-1]
    base = None
    if tag == "srgbClr" and el.get("val"):
        base = hex_to_rgb(el.get("val"))
    elif tag == "schemeClr" and el.get("val"):
        name = SCHEME_ALIAS.get(el.get("val"), el.get("val"))
        if name in theme:
            base = hex_to_rgb(theme[name])
    elif tag == "sysClr" and el.get("lastClr"):
        base = hex_to_rgb(el.get("lastClr"))
    elif tag == "prstClr" and el.get("val"):
        base = PRESET_COLOURS.get(el.get("val").lower())
    if base is None:
        return None, True
    return _apply_transforms(base, el)


COLOUR_TAGS = ("srgbClr", "schemeClr", "sysClr", "prstClr", "hslClr", "scrgbClr")

PRESET_COLOURS = {"black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
                  "green": (0, 128, 0), "blue": (0, 0, 255), "yellow": (255, 255, 0),
                  "gray": (128, 128, 128), "grey": (128, 128, 128),
                  "darkGray".lower(): (169, 169, 169), "lightgray": (211, 211, 211)}


def _first_colour_child(parent):
    if parent is None:
        return None
    for child in parent:
        if child.tag.split("}")[-1] in COLOUR_TAGS:
            return child
    return None


def _solid_to_rgb(fill, theme):
    """An <a:solidFill> element -> concrete RGB, or UNKNOWN_FILL if unmodellable."""
    if fill is None:
        return None
    rgb, ok = _colour_el_to_rgb(_first_colour_child(fill), theme)
    if rgb is None:
        return None
    return rgb if ok else UNKNOWN_FILL


def theme_fill_styles(prs):
    """The theme's fillStyleLst, which <a:fillRef idx="n"> points into (1-based)."""
    try:
        part = prs.slide_master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        root = getattr(part, "_element", None)
        if root is None:
            from lxml import etree
            root = etree.fromstring(part.blob)
        lst = root.find(f".//{A}fmtScheme/{A}fillStyleLst")
        return list(lst) if lst is not None else []
    except Exception:
        return []


def _fill_from_style_ref(shape, theme, fill_styles):
    """Resolve <p:style><a:fillRef>, which is how MOST real shapes are filled.

    A shape drawn in PowerPoint, and every shape python-pptx's add_shape() emits,
    carries NO fill in spPr at all -- the fill comes from this style reference into
    the theme's fillStyleLst, with a colour override on the ref itself. Reading only
    spPr therefore reports the single most common shape in any deck as "no fill",
    falls through to the slide background, and fabricates BOTH sides of the contrast
    pair. Verified: python-pptx add_shape writes fillRef idx=3 schemeClr accent1 and
    nothing in spPr.
    """
    try:
        ref = shape._element.find(f".//{P}style/{A}fillRef")
        if ref is None:
            return None
        idx = int(ref.get("idx") or 0)
        if idx == 0:
            return None                      # idx 0 is explicitly no fill
        rgb, ok = _colour_el_to_rgb(_first_colour_child(ref), theme)
        if rgb is None:
            return None
        if not ok:
            return UNKNOWN_FILL
        # The referenced style may be a gradient or pattern, in which case the
        # override colour is only one stop of it and a flat answer would be wrong.
        style = fill_styles[idx - 1] if 0 < idx <= len(fill_styles) else None
        if style is not None and style.tag.split("}")[-1] != "solidFill":
            return UNKNOWN_FILL
        return rgb
    except Exception:
        return None


def _fill_of(props, theme):
    """Resolve the fill on a properties element (spPr / bgPr / grpSpPr).

    Only DIRECT children are considered. A descendant search here is a real bug:
    `.//solidFill` on a shape walks into <a:txBody> and returns the TEXT colour as
    the shape fill, which makes foreground and background identical and reports a
    perfect 1.0:1 on a shape that has no fill at all. It also means the true slide
    background is never used, so genuine contrast failures go unmeasured.
    <a:ln> (the outline colour) is a direct child too, and is likewise not the fill.
    """
    if props is None:
        return None
    for child in props:
        tag = child.tag
        if tag == f"{A}solidFill":
            return _solid_to_rgb(child, theme)
        if tag == f"{A}noFill":
            return None                 # explicitly transparent -> inherit
        if tag in (f"{A}blipFill", f"{A}gradFill", f"{A}pattFill"):
            return UNKNOWN_FILL
    return None                          # no fill specified -> inherit


def shape_fill_rgb(shape, theme, fill_styles=()):
    """Concrete RGB for a shape's backdrop, or UNKNOWN_FILL, or None to inherit.

    Order matters: spPr wins, then the p:style fill reference. A picture or a
    style-driven table is UNKNOWN, never a silent fall-through to the slide
    background -- that fall-through is what let this harness report a fabricated
    1.0:1 on white-over-photo and a fabricated PASS on black-over-black-photo.
    """
    try:
        el = shape._element
        tag = el.tag.split("}")[-1]

        # <p:pic> keeps its image in blipFill as a DIRECT CHILD, never inside spPr.
        if tag == "pic" or el.find(f"{A}blipFill") is not None or el.find(f"{P}blipFill") is not None:
            return UNKNOWN_FILL

        for t in (f"{P}spPr", f"{A}spPr", f"{P}grpSpPr", f"{A}grpSpPr"):
            props = el.find(t)
            if props is not None:
                got = _fill_of(props, theme)
                if got is not None:
                    return got
                break

        got = _fill_from_style_ref(shape, theme, fill_styles)
        if got is not None:
            return got

        # A table's banding lives in tableStyles.xml, which is a whole style system
        # this does not read. Guessing white here fabricated nine contrast failures
        # on a default-styled table over a dark panel.
        if getattr(shape, "has_table", False) and shape.has_table:
            return UNKNOWN_FILL
    except Exception:
        pass
    return None


def run_colour_rgb(run, theme):
    """A run's OWN explicit colour, or None. Transforms are applied; an unmodellable
    one returns None so the caller reports it unmeasured rather than guessing."""
    try:
        rpr = run._r.find(f"{A}rPr")
        if rpr is not None:
            fill = rpr.find(f"{A}solidFill")
            if fill is not None:
                rgb, ok = _colour_el_to_rgb(_first_colour_child(fill), theme)
                if rgb is not None:
                    return rgb if ok else None
    except Exception:
        pass
    return None


def slide_bg_rgb(slide, prs, theme):
    """Background for a slide: its own, else the layout's, else the master's.

    Returns UNKNOWN_FILL for a picture or gradient background rather than falling
    through to a white guess. Text over a photo cannot be contrast-checked from the
    file, and pretending otherwise is worse than saying so.
    """
    for src in (slide, slide.slide_layout, prs.slide_master):
        try:
            bg = src.element.find(f".//{P}bg")
            if bg is None:
                continue
            props = bg.find(f"{P}bgPr")
            if props is not None:
                got = _fill_of(props, theme)
                if got is not None:
                    return got
            # <p:bgRef idx="..."> points at the theme's fill style list and carries
            # its own colour override, which is the usual shape of a branded master.
            ref = bg.find(f"{P}bgRef")
            if ref is not None:
                got = _solid_to_rgb(ref, theme)
                if got:
                    return got
        except Exception:
            continue
    return (255, 255, 255)


# ── traversal ────────────────────────────────────────────────────────────────

def _group_transform(grp):
    """Map a group's child coordinate space onto the slide.

    A group has an offset/extent on the slide (off/ext) AND a child coordinate
    system (chOff/chExt) which is usually the same numbers but need not be. A child
    at chOff sits at off, and the space scales by ext/chExt. Ignoring this reports
    grouped shapes at their raw child coordinates, which are meaningless on the slide.
    """
    try:
        # DIRECT child of grpSpPr, never a descendant search. `.//xfrm` on a group
        # with a bare <p:grpSpPr/> (legal, and several exporters emit it) returns the
        # FIRST CHILD SHAPE's xfrm, translating every child by that child's own
        # offset. Third instance of this exact bug class in this file.
        props = _find_any(grp._element, f"{P}grpSpPr", f"{A}grpSpPr")
        xfrm = props.find(f"{A}xfrm") if props is not None else None
        if xfrm is None:
            return (0, 0, 1.0, 1.0, 0, 0)
        off, ext = xfrm.find(f"{A}off"), xfrm.find(f"{A}ext")
        cho, che = xfrm.find(f"{A}chOff"), xfrm.find(f"{A}chExt")
        ox, oy = int(off.get("x")), int(off.get("y"))
        cx, cy = (int(cho.get("x")), int(cho.get("y"))) if cho is not None else (0, 0)
        sx = (int(ext.get("cx")) / int(che.get("cx"))) if (ext is not None and che is not None
                                                          and int(che.get("cx"))) else 1.0
        sy = (int(ext.get("cy")) / int(che.get("cy"))) if (ext is not None and che is not None
                                                          and int(che.get("cy"))) else 1.0
        return (ox, oy, sx, sy, cx, cy)
    except Exception:
        return (0, 0, 1.0, 1.0, 0, 0)


def _compose(outer, inner):
    """Nest one group transform inside another."""
    ox, oy, sx, sy, cx, cy = outer
    ix, iy, isx, isy, icx, icy = inner
    return (ox + (ix - cx) * sx, oy + (iy - cy) * sy, sx * isx, sy * isy, icx, icy)


def _apply(tf, x, y):
    ox, oy, sx, sy, cx, cy = tf
    return ox + (x - cx) * sx, oy + (y - cy) * sy


IDENTITY = (0, 0, 1.0, 1.0, 0, 0)


def walk_shapes(shapes, tf=IDENTITY, depth=0):
    """Yield every (shape, transform, in_group) on a slide, recursing into GROUPS.

    🔴 Without this the harness is blind to exactly the decks it exists to check.
    deck-forge MANDATES building every visual from grouped rectangles and text boxes
    (because native charts flatten on Google Slides import), so a non-recursive
    `for shape in slide.shapes` measures none of the real content. Tested: a slide of
    grouped 7pt bars plus a 6pt table cell reported "0 runs, all PASS, exit 0".
    """
    if depth > 12:
        return
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield shape, tf, depth > 0
            child_tf = _compose(tf, _group_transform(shape))
            yield from walk_shapes(shape.shapes, child_tf, depth + 1)
        else:
            yield shape, tf, depth > 0


def text_frames_of(shape):
    """Every text frame a shape carries, INCLUDING every table cell.

    A table is one shape with no `.text_frame`; its text lives in cell frames. Skipping
    them silently drops every number in every table, which on a data deck is the
    majority of the text on the slide.
    """
    out = []
    if getattr(shape, "has_table", False) and shape.has_table:
        for r, row in enumerate(shape.table.rows):
            for c, cell in enumerate(row.cells):
                out.append((cell.text_frame, f"table[{r}][{c}]", cell))
    elif shape.has_text_frame:
        out.append((shape.text_frame, None, None))
    return out


def _rot_deg(shape):
    """@rot on the shape's xfrm, in degrees. OOXML stores 60000ths of a degree."""
    try:
        props = _find_any(shape._element, f"{P}spPr", f"{A}spPr")
        xfrm = props.find(f"{A}xfrm") if props is not None else None
        if xfrm is not None and xfrm.get("rot"):
            return (int(xfrm.get("rot")) / 60000.0) % 360.0
    except Exception:
        pass
    return 0.0


def _bounds(shape, tf):
    """Absolute slide-coordinate bounds in EMU, or None.

    Accounts for @rot, which PowerPoint applies ABOUT THE SHAPE CENTRE. Ignoring it
    is wrong in both directions: an 8x1in bar rotated 90 degrees renders 1x8in and
    hangs off a 7.5in slide while the raw box fits, and a 1x8in box rotated 90
    renders 8x1in fully on-slide while the raw box reports off-slide. It also
    corrupts _contains, so a rotated panel gets credited as a backdrop it is not.
    """
    try:
        if shape.left is None or shape.top is None:
            return None
        w, h = (shape.width or 0), (shape.height or 0)
        rot = _rot_deg(shape)
        if rot:
            cx, cy = shape.left + w / 2.0, shape.top + h / 2.0
            rad = math.radians(rot)
            ac, as_ = abs(math.cos(rad)), abs(math.sin(rad))
            # Axis-aligned bounding box of the rotated rectangle.
            bw, bh = w * ac + h * as_, w * as_ + h * ac
            corners = [(cx - bw / 2.0, cy - bh / 2.0), (cx + bw / 2.0, cy + bh / 2.0)]
        else:
            corners = [(shape.left, shape.top), (shape.left + w, shape.top + h)]
        pts = [_apply(tf, x, y) for x, y in corners]
        xs, ys = [q[0] for q in pts], [q[1] for q in pts]
        return (min(xs), min(ys), max(xs), max(ys))
    except Exception:
        return None


def _contains(outer, inner, slack=9144):    # 0.01in of slack
    if not outer or not inner:
        return False
    return (outer[0] <= inner[0] + slack and outer[1] <= inner[1] + slack
            and outer[2] >= inner[2] - slack and outer[3] >= inner[3] - slack)


def backdrop_for(i, layers, slide_bg):
    """What is actually BEHIND layers[i]'s text.

    A shape with no fill of its own is not sitting on the slide background if some
    filled shape is painted underneath it. Placing a text box over a coloured
    rectangle is one of the most common ways a real deck is built, and without this
    the harness reports white-on-white at 1.0:1 for perfectly legible text -- a false
    positive, which is the failure that makes a checker untrustworthy.

    So: walk DOWN the z-order from just below this shape and take the first filled
    shape that fully contains it.
    """
    own = layers[i]["fill"]
    if own is not None:
        return own, "own fill"
    me = layers[i]["bounds"]
    for j in range(i - 1, -1, -1):
        under = layers[j]
        if under["fill"] is None or not _contains(under["bounds"], me):
            continue
        if under["fill"] is UNKNOWN_FILL:
            return UNKNOWN_FILL, "sits on a picture or gradient"
        return under["fill"], f"behind it: {under['name']}"
    return slide_bg, "slide background"


# ── the checks ───────────────────────────────────────────────────────────────

# ── the legibility model, and exactly what it is drawn from ─────────────────
#
# ANSI/INFOCOMM V202.01:2016 §4.3.1:  IH = FV / (200 x %EH)
#   "200 is the Acuity Factor for Basic Decision Making", derived in §6.1.2 from a
#   minimum resolvable angle of 17.25 arcminutes. Rearranged for absolute element
#   height EH = %EH x IH, that is FV = 200 x EH -- the "element height >= viewing
#   distance / 200" form used here.
#
# 🔴 THE ELEMENT FOR TEXT IS A LOWERCASE LETTER, NOT A CAPITAL.
#   AVIXA's own CTS-Prep material, slide 27 "What is Percent Element Height?":
#       "With text: lowercase letter"
#   (identical wording in the 2017 Italy and 2019 France decks). The published
#   standard's §3.2.6 is ambiguous; AVIXA's teaching material is not. This harness
#   originally measured CAP height, which made every floor about 35% too low --
#   lenient, which is the dangerous direction for a legibility check.
#
# ⚠️ THE 4/6/8 RULE IS NOT PART OF THE STANDARD. It appears zero times in V202.01,
#   and AVIXA's CTS-Prep deck puts it on a slide titled "The Old Way of Doing Things
#   4/6/8", listing "Only a Best Practice" and "Origins are unclear". The standard's
#   own Foreword says prior methods are "not attributable to any particular source
#   and appear to be based on precedent". It survives here only because it is
#   numerically identical to the standards-backed part: FV = 200 x %EH x IH makes
#   2/3/4 %EH give exactly 4/6/8 screen heights. So the ARITHMETIC is sound and the
#   PROVENANCE is folklore. Do not describe this model as "not folklore".
#
# ⚠️ The tiers below are all Basic Decision Making at different %EH. DISCAS
#   Analytical Decision Making is a DIFFERENT calculation, IH = (IR x FV) / 3438,
#   driven by vertical pixel resolution with no %EH term. "analytical" here means
#   2%EH, not DISCAS ADM.
#
# AVIXA CTS-Prep slide 39: "Typically use something from 2% Element Height to 4%
# Element Height ... A 3% Element Height is a good starting point."

PERCENT_EH = {"analytical": 0.02, "basic": 0.03, "passive": 0.04}
VIEWING_NEED = {"analytical": 4.0, "basic": 6.0, "passive": 8.0}   # screen heights

# x-height / em, measured from OS/2.sxHeight / head.unitsPerEm on the real binaries.
X_HEIGHT_RATIO = {
    "arial": 0.5186, "helvetica": 0.5229, "helvetica neue": 0.5170,
    "verdana": 0.5454, "tahoma": 0.5460, "calibri": 0.4775, "carlito": 0.4775,
    "georgia": 0.4814, "times new roman": 0.4473, "times": 0.4473,
    "garamond": 0.4400, "inter": 0.7275 * 0.72, "roboto": 0.5283,
    "open sans": 0.5350, "lato": 0.5060, "segoe ui": 0.5000,
}
DEFAULT_X_RATIO = 0.52       # Arial/Helvetica class, the common deck faces
CAP_HEIGHT_RATIO = 0.70      # kept for --element cap
SUBTENSE_DIVISOR = 200.0     # the Basic Decision Making acuity factor


def glyph_ratio(font_name=None, element="x"):
    """Height of the reference glyph as a fraction of the point size."""
    if element == "cap":
        return CAP_HEIGHT_RATIO, "cap height"
    if font_name:
        r = X_HEIGHT_RATIO.get(str(font_name).strip().lower())
        if r:
            return r, f"x-height of {font_name}"
    return DEFAULT_X_RATIO, "x-height (generic)"


def legibility_floor_pt(slide_h_in, need="basic", screen_h_in=None, room_depth_ft=None,
                        font_name=None, element="x"):
    """Minimum legible point size on a projected slide.

    A point is a DOCUMENT unit, not a physical one. Its height on the wall depends
    entirely on how large the slide is projected:

        fraction_of_image = (pt / 72) * glyph_ratio / slide_h_in
        physical_glyph_in = fraction_of_image * screen_h_in
        require             physical_glyph_in >= viewing_distance_in / 200

    Treating points as physical inches reports a ~185pt floor for a 30ft room.

    With no measured screen, screen height comes from the %EH tier, and the viewing
    distance then CANCELS -- not a discovered invariant, just the substitution, since
    screen height is defined as distance/ratio. Consequence: --room-depth alone can
    never turn a pass into a fail. Pass --screen-height to measure a real room.

    ⚠️ The result is a property of SLIDE HEIGHT IN INCHES, not of aspect ratio.
    16:9 at 13.333x7.5 and 16:9 at 10x5.625 are the same shape and give 31.2pt vs
    23.4pt, because the second authors everything at 75% scale.
    """
    ratio, _ = glyph_ratio(font_name, element)
    if screen_h_in and room_depth_ft:
        need_in = (room_depth_ft * 12.0) / SUBTENSE_DIVISOR
        return round(need_in / screen_h_in * slide_h_in * 72.0 / ratio, 1)
    pct = PERCENT_EH.get(need, PERCENT_EH["basic"])
    return round(pct * slide_h_in * 72.0 / ratio, 1)


def audit(path, room_depth_ft=None, min_pt=None, need="basic", screen_h_in=None,
          element="x", font_name=None):
    prs = Presentation(str(path))
    theme = theme_colours(prs)
    mstyles = master_txstyles(prs)
    fill_styles = theme_fill_styles(prs)
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN

    derived_pt = legibility_floor_pt(slide_h_in, need, screen_h_in, room_depth_ft,
                                     font_name, element)
    _, glyph_desc = glyph_ratio(font_name, element)
    # What screen the stated room actually needs, which is the useful thing the
    # room depth tells you once the point-size floor stops depending on it.
    needed_screen_h_ft = (room_depth_ft / VIEWING_NEED.get(need, 6.0)) if room_depth_ft else None

    floor_pt = min_pt or derived_pt

    findings = {"tiny_text": [], "native_charts": [], "low_contrast": [], "offslide": [],
                "unresolved_size": [], "empty_placeholders": [], "unmeasured_contrast": []}
    sizes, checked_runs, checked_shapes = set(), 0, 0
    per_slide = []

    for n, slide in enumerate(prs.slides, 1):
        bg = slide_bg_rgb(slide, prs, theme)
        s_words, s_shapes = 0, 0
        ordered = list(walk_shapes(slide.shapes))
        layers = [{"fill": shape_fill_rgb(sh, theme, fill_styles), "bounds": _bounds(sh, t),
                   "name": sh.name} for sh, t, _ in ordered]
        for zi, (shape, tf, in_group) in enumerate(ordered):
            checked_shapes += 1
            s_shapes += 1
            loc = f"slide {n}" + (" (in group)" if in_group else "")

            # native chart -> flattens to a dead image on Google Slides import
            if getattr(shape, "has_chart", False) and shape.has_chart:
                findings["native_charts"].append({
                    "where": loc, "name": shape.name,
                    "why": "a native chart converts to a STATIC IMAGE on Google Slides import and "
                           "cannot be made editable — Slides only keeps a chart live when linked to "
                           "a Sheet, and that link cannot exist in an uploaded .pptx. Rebuild from "
                           "rectangles and text boxes."})

            # off-slide / bleeding geometry, in SLIDE coordinates
            try:
                bb = layers[zi]["bounds"]
                if bb is not None:
                    l, t = bb[0] / EMU_PER_IN, bb[1] / EMU_PER_IN
                    r, b = bb[2] / EMU_PER_IN, bb[3] / EMU_PER_IN
                    if l < -0.05 or t < -0.05 or r > slide_w_in + 0.05 or b > slide_h_in + 0.05:
                        findings["offslide"].append({
                            "where": loc, "name": shape.name,
                            "detail": f"({l:.2f},{t:.2f})–({r:.2f},{b:.2f})in vs "
                                      f"{slide_w_in:.2f}×{slide_h_in:.2f}in"})
            except Exception:
                pass

            frames = text_frames_of(shape)
            if not frames:
                continue
            shape_fill, fill_src = backdrop_for(zi, layers, bg)

            for tframe, cell_ref, cell in frames:
                where = loc + (f" {cell_ref}" if cell_ref else "")
                fill = shape_fill
                if cell is not None:
                    cf = _fill_of(cell._tc.find(f"{A}tcPr"), theme)
                    if cf is not None:
                        fill = cf
                text_here = tframe.text.strip()
                if not text_here:
                    idx, _ = _ph_idx(shape)
                    if idx is not None and cell is None:
                        findings["empty_placeholders"].append({"where": where, "name": shape.name})
                    continue
                s_words += len(text_here.split())

                for para in tframe.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        checked_runs += 1
                        pt, src = resolve_font_pt(run, para, shape, slide, prs, mstyles)
                        sizes.add(round(pt, 1))
                        if src.startswith("PowerPoint default"):
                            findings["unresolved_size"].append({
                                "where": where, "text": run.text.strip()[:40],
                                "why": "no size found anywhere in the chain; assumed 18pt"})
                        if pt < floor_pt:
                            findings["tiny_text"].append({
                                "where": where, "pt": pt, "floor": floor_pt, "source": src,
                                "text": run.text.strip()[:48]})
                        fg, fg_src = resolve_font_colour(
                            run, para, shape, slide, prs, mstyles, theme)
                        if fg is None:
                            # Inherited from theme/placeholder. Say so; do not assume black.
                            findings["unmeasured_contrast"].append({
                                "where": where, "why": "text colour resolves nowhere in the chain",
                                "text": run.text.strip()[:40]})
                        elif fill is UNKNOWN_FILL:
                            findings["unmeasured_contrast"].append({
                                "where": where, "why": f"background is indeterminate ({fill_src})",
                                "text": run.text.strip()[:40]})
                        else:
                            ratio = contrast(fg, fill)
                            large = is_large_text(pt, run.font.bold)
                            need_ratio = 3.0 if large else 4.5
                            if ratio < need_ratio:
                                findings["low_contrast"].append({
                                    "where": where, "ratio": round(ratio, 2), "need": need_ratio,
                                    "pt": pt, "text": run.text.strip()[:40]})
        per_slide.append({"slide": n, "shapes": s_shapes, "words": s_words})

    return {
        "file": Path(path).name,
        "slides": len(prs.slides),
        "aspect": f"{slide_w_in:.2f}x{slide_h_in:.2f}in",
        "checked": {"shapes": checked_shapes, "runs": checked_runs},
        "type_sizes": sorted(sizes),
        "legibility_floor_pt": floor_pt,
        "legibility_source": ("explicit --min-pt" if min_pt else
                             f"{need} viewing need ({PERCENT_EH.get(need, 0.03) * 100:g}%EH), "
                             f"{glyph_desc}"),
        "room_depth_ft": room_depth_ft,
        "needed_screen_height_ft": round(needed_screen_h_ft, 1) if needed_screen_h_ft else None,
        "findings": findings,
        "per_slide": per_slide,
        "theme_colours": theme,
    }


def exit_code(r, fail_on):
    """The ONE place a verdict is computed.

    --json used to recompute this locally and skipped both the zero-run trap and
    --fail-on any, so the mode every CI gate and agent actually uses reported 0 on a
    deck with no measurable text -- precisely the state this module calls the most
    dangerous it can be in.
    """
    n = {k: len(v) for k, v in r["findings"].items()}
    hard = n["native_charts"] + n["tiny_text"] + n["low_contrast"] + n["offslide"]
    total = sum(n.values())
    if r["checked"]["runs"] == 0:
        hard += 1
        total += 1
    if fail_on == "none":
        return 0
    return 1 if (total if fail_on == "any" else hard) else 0


def report(r, fail_on):
    f = r["findings"]
    n = {k: len(v) for k, v in f.items()}
    ok = lambda c: "PASS" if c == 0 else "FAIL"
    print(f"\n{r['file']} — {r['slides']} slides, {r['aspect']}")
    print(f"  checked {r['checked']['shapes']} shapes / {r['checked']['runs']} text runs")
    if r["checked"]["runs"] == 0:
        print("  🔴 NO TEXT RUNS MEASURED — not a pass. The deck may be image-only, "
              "or the reader failed.")
    print()
    print(f"  nativeCharts   {n['native_charts']} {ok(n['native_charts'])}"
          f"   — flatten to dead images on Google Slides import")
    print(f"  legibility     {n['tiny_text']} below {r['legibility_floor_pt']}pt "
          f"{ok(n['tiny_text'])}   — {r['legibility_source']}")
    print(f"  contrast       {n['low_contrast']} WCAG AA failures {ok(n['low_contrast'])}")
    print(f"  offSlide       {n['offslide']} {ok(n['offslide'])}")
    print(f"  emptyPlaceholder {n['empty_placeholders']} {ok(n['empty_placeholders'])}")
    print(f"  unresolvedSize {n['unresolved_size']}"
          + ("  ⚠️  assumed 18pt — verify by hand" if n["unresolved_size"] else ""))
    print(f"  unmeasured     {n['unmeasured_contrast']} run(s) whose contrast could NOT be "
          f"computed" + ("  ⚠️  not a pass" if n["unmeasured_contrast"] else ""))
    print(f"  typeSizes      {r['type_sizes']}")
    if r["needed_screen_height_ft"]:
        print(f"\n  A {r['room_depth_ft']:.0f}ft room at the '{r['legibility_source']}' needs a "
              f"screen at least {r['needed_screen_height_ft']}ft high")
        print(f"  ({r['needed_screen_height_ft'] * 16 / 9:.1f}ft wide at 16:9). Undersize the "
              f"screen and no point size saves you.")
    print()
    for key, label, mark in (("native_charts", "NATIVE CHART", "🔴"),
                             ("tiny_text", "TOO SMALL", "🔴"),
                             ("low_contrast", "LOW CONTRAST", "⚠️ "),
                             ("offslide", "OFF-SLIDE", "⚠️ ")):
        for item in f[key][:8]:
            detail = item.get("why") or item.get("detail") or ""
            extra = ""
            if key == "tiny_text":
                extra = f" {item['pt']}pt (from {item['source']}) “{item['text']}”"
            elif key == "low_contrast":
                extra = f" {item['ratio']}:1 need {item['need']} “{item['text']}”"
            print(f"  {mark} [{label}] {item['where']}{extra}")
            if detail:
                print(f"      {detail}")
    print("\n  caveat: static file analysis. Says nothing about whether the deck is any good,")
    print("  whether the story works, or how it reads when projected in a lit room.\n")
    return exit_code(r, fail_on)


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("file", type=Path)
    ap.add_argument("--room-depth", type=float, default=None,
                    help="farthest viewer distance in feet; reports the screen size it demands")
    ap.add_argument("--screen-height", type=float, default=None,
                    help="actual screen height in INCHES; with --room-depth this measures the "
                         "real room instead of assuming the 4/6/8 rule")
    ap.add_argument("--viewing-need", choices=list(VIEWING_NEED), default="basic",
                    help="analytical (dense tables) / basic (default) / passive (back of a hall)")
    ap.add_argument("--element", choices=["x", "cap"], default="x",
                    help="reference glyph. AVIXA's own training says the element for "
                         "text is a LOWERCASE letter, so x-height is the default; "
                         "'cap' is ~35%% more lenient and not what the standard teaches")
    ap.add_argument("--font", default=None,
                    help="font name, for a real x-height ratio instead of the generic 0.52")
    ap.add_argument("--min-pt", type=float, default=None, help="override the point-size floor")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--fail-on", choices=["hard", "any", "none"], default="hard")
    a = ap.parse_args()
    if not a.file.exists():
        print(f"ERROR: file not found: {a.file}", file=sys.stderr)
        return 1
    if a.screen_height and not a.room_depth:
        print("ERROR: --screen-height needs --room-depth to mean anything.", file=sys.stderr)
        return 2
    try:
        r = audit(a.file, a.room_depth, a.min_pt, a.viewing_need, a.screen_height,
                  a.element, a.font)
    except Exception as e:
        print(f"ERROR: could not read {a.file}: {type(e).__name__}: {e}", file=sys.stderr)
        return 1
    if a.json:
        print(json.dumps(r, indent=2))
        return exit_code(r, a.fail_on)
    return report(r, a.fail_on)


if __name__ == "__main__":
    sys.exit(main())
